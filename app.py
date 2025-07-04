from flask import render_template, request, redirect, url_for, session, flash
import pyodbc
from config import app, DB_CONFIG
import os
import traceback
import tempfile
import docx2txt
import pptx
import PyPDF2
import openai
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

def extract_text(file_path, filename):
    ext = filename.lower().split('.')[-1]
    text = ""
    if ext == "pdf":
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
    elif ext in ["doc", "docx"]:
        text = docx2txt.process(file_path)
    elif ext in ["ppt", "pptx"]:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    return text.strip()


# DB connection
def get_db_connection():
    return pyodbc.connect(
        f"DRIVER={{SQL Server}};SERVER={DB_CONFIG['server']};DATABASE={DB_CONFIG['database']};UID={DB_CONFIG['username']};PWD={DB_CONFIG['password']}"
    )

# Home page (redirected to after login/register)
@app.route('/')
def index():
    if 'username' in session:
        return render_template('home.html', username=session['username'])
    return redirect(url_for('login'))

# LOGIN
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        identifier = request.form['identifier']
        password = request.form['password']

        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT username FROM users 
            WHERE (username = ? OR email = ?) AND password = ?
        """, (identifier, identifier, password))
        result = cursor.fetchone()
        conn.close()

        if result:
            session['username'] = result[0]
            return redirect(url_for('home'))
        else:
            flash('Invalid login credentials.')
    return render_template('login.html')

# REGISTER
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        dob = request.form['dob']

        conn = get_db_connection()
        cursor = conn.cursor()

        # Check uniqueness
        cursor.execute("SELECT 1 FROM users WHERE username = ? OR email = ?", (username, email))
        if cursor.fetchone():
            conn.close()
            flash('Username or email already exists.')
            return render_template('register.html')

        # Insert user
        cursor.execute("""
            INSERT INTO users (username, email, password, date_of_birth)
            VALUES (?, ?, ?, ?)
        """, (username, email, password, dob))
        conn.commit()
        conn.close()

        session['username'] = username
        return redirect(url_for('home'))

    return render_template('register.html')

@app.route('/logout')
def logout():
    session.pop('username', None)
    return redirect(url_for('login'))



@app.route('/account', methods=['GET', 'POST'])
def account():
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    if request.method == 'POST':
        new_email = request.form['email'].strip()
        new_password = request.form['password'].strip()
        dob = request.form['dob']
        username = session['username']

        # Check if email is already taken by another user
        cursor.execute("""
            SELECT * FROM users
            WHERE email = ? AND username != ?
        """, (new_email, username))
        conflict = cursor.fetchone()

        if conflict:
            cursor.execute("SELECT * FROM users WHERE username = ?", (username,))
            user = cursor.fetchone()
            conn.close()
            return render_template("account.html", user={
                'username': user[0], 'email': user[1], 'password': user[2],
                'dob': user[3], 'quiz_completed': user[4]
            }, error="Email already in use.")

        # ✅ Update only email, password, dob
        cursor.execute("""
            UPDATE users
            SET email = ?, password = ?, date_of_birth = ?
            WHERE username = ?
        """, (new_email, new_password, dob, username))

        conn.commit()
        conn.close()

        return redirect(url_for('account'))

    # GET request – Load user info
    cursor.execute("SELECT * FROM users WHERE username = ?", (session['username'],))
    user = cursor.fetchone()
    conn.close()

    return render_template("account.html", user={
        'username': user[0], 'email': user[1], 'password': user[2],
        'dob': user[3], 'quiz_completed': user[4]
    })

@app.route('/home', methods=['GET', 'POST'])
def home():
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()
    success = error = None

    if request.method == 'POST':
        course_name = request.form['course_name'].strip()
        username = session['username']

        cursor.execute("SELECT * FROM courses WHERE course_name = ? AND username = ?", (course_name, username))
        if cursor.fetchone():
            error = "Course name already exists."
        else:
            cursor.execute("INSERT INTO courses (course_name, username) VALUES (?, ?)", (course_name, username))
            conn.commit()
            success = f"Course '{course_name}' added successfully."

    cursor.execute("SELECT * FROM courses WHERE username = ?", (session['username'],))
    courses = cursor.fetchall()
    conn.close()
    return render_template("home.html", courses=courses, error=error, success=success)

import traceback

@app.route('/course/<int:course_id>', methods=['GET', 'POST'])
def view_course(course_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    # Pop flash messages first
    error = session.pop('chapter_error', None)
    success = session.pop('chapter_success', None)

    if request.method == 'POST':
        print("📥 POST /course/<course_id> triggered.")
        chapter_name = request.form['chapter_name'].strip()
        file = request.files.get('chapter_file')
        error = None
        conn = None

        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            print("🔌 DB connected for POST")

            # Verify course belongs to user
            cursor.execute("SELECT course_name FROM courses WHERE course_id = ? AND username = ?", 
                           (course_id, session['username']))
            course = cursor.fetchone()
            if not course:
                print("❌ Invalid course or permission denied")
                conn.close()
                return redirect(url_for('home'))

            # Validate file
            if not file:
                error = "No file uploaded."
            elif file.filename.split('.')[-1].lower() not in ['pdf', 'doc', 'docx', 'ppt', 'pptx', 'txt']:
                error = "Unsupported file type."
            else:
                file.seek(0)
                if len(file.read()) > 5 * 1024 * 1024:
                    error = "File too large. Limit is 5MB."
                else:
                    file.seek(0)
                    with tempfile.NamedTemporaryFile(delete=False) as tmp:
                        file.save(tmp.name)
                        tmp_path = tmp.name
                        print("📄 File saved at", tmp_path)
                    text = extract_text(tmp_path, file.filename)
                    os.unlink(tmp_path)
                    print("📄 Extracted text length:", len(text) if text else "None")

                    if not text:
                        error = "Failed to extract text from file."
                    else:
                        prompt = f"Summarize the following file content in a friendly, easy-to-understand way:\n{text[:4000]}"
                        print("🤖 Sending to OpenAI")
                        try:
                            client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {"role": "system", "content": "You are a helpful assistant who explains things clearly."},
                                    {"role": "user", "content": prompt}
                                ]
                            )
                            summary = response.choices[0].message.content.strip()
                            print("✅ Got summary")

                            cursor.execute(
                                "INSERT INTO chapters (chapter_name, course_id, chapter_summary) VALUES (?, ?, ?)",
                                (chapter_name, course_id, summary)
                            )
                            conn.commit()
                            print("✅ Inserted chapter and committed")
                            success = f"Chapter '{chapter_name}' added successfully."

                        except Exception as e:
                            error = f"AI summary failed: {str(e)}"
                            print("❌ OpenAI or DB error:")
                            traceback.print_exc()

            if error:
                session['chapter_error'] = error
            else:
                session['chapter_success'] = success

        except Exception as e:
            print("❌ Fatal POST error:")
            traceback.print_exc()
            error = "Internal error while adding chapter."
            session['chapter_error'] = error
        finally:
            if conn:
                conn.close()
                print("🔒 DB connection closed (POST)")

        return redirect(url_for('view_course', course_id=course_id))

    # GET request
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        print("🌐 GET request: Fetching course and chapters")
        cursor.execute("SELECT course_name FROM courses WHERE course_id = ? AND username = ?", 
                       (course_id, session['username']))
        course = cursor.fetchone()
        if not course:
            conn.close()
            return redirect(url_for('home'))
        course_name = course[0]

        cursor.execute("SELECT * FROM chapters WHERE course_id = ?", (course_id,))
        chapters = cursor.fetchall()
        conn.close()
        print("✅ GET complete")
    except Exception as e:
        print("❌ GET error:")
        traceback.print_exc()
        course_name = "Unknown"
        chapters = []
        error = "Could not load course or chapters."

    return render_template("chapter.html",
                           course_name=course_name,
                           chapters=chapters,
                           error=error,
                           success=success)

@app.route('/summary/<int:chapter_id>')
def view_summary(chapter_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get chapter
    cursor.execute("""
        SELECT c.chapter_name, c.chapter_summary, cr.course_name
        FROM chapters c
        JOIN courses cr ON c.course_id = cr.course_id
        WHERE c.chapter_id = ?
    """, (chapter_id,))
    chapter = cursor.fetchone()
    conn.close()

    if not chapter:
        return redirect(url_for('home'))

    return render_template("summary.html",
                           chapter_name=chapter.chapter_name,
                           summary=chapter.chapter_summary,
                           course_name=chapter.course_name,
                           chapter_id=chapter_id)

@app.route('/generate_quiz/<int:chapter_id>', methods=['POST'])
def generate_quiz(chapter_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    try:
        conn = get_db_connection()
        cursor = conn.cursor()

        # Get chapter summary
        cursor.execute("SELECT chapter_summary FROM chapters WHERE chapter_id = ?", (chapter_id,))
        row = cursor.fetchone()
        if not row:
            conn.close()
            return "Chapter not found.", 404

        summary = row.chapter_summary

        # Prepare OpenAI prompt
        prompt = (
            f"Based on the following summary, generate a 5-question multiple choice quiz. "
            f"Each question should have exactly 3 options (a, b, c) and only one correct answer. "
            f"Format:\n"
            f"QUESTION:\n"
            f"1) Question text\n"
            f"a) Option A\n"
            f"b) Option B\n"
            f"c) Option C\n"
            f"2) ...\n"
            f"ANSWER:\n"
            f"1) a\n2) b\n..."
            f"\n\nSUMMARY:\n{summary[:4000]}"
        )

        client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant who creates quizzes."},
                {"role": "user", "content": prompt}
            ]
        )

        full_text = response.choices[0].message.content.strip()

        # Separate QUESTION and ANSWER
        if "ANSWER:" in full_text:
            quiz_part, answer_part = full_text.split("ANSWER:", 1)
        else:
            quiz_part = full_text
            answer_part = "Could not extract answers."

        # Insert quiz into database
        cursor.execute(
            "INSERT INTO quizzes (chapter_id, quiz_content, quiz_answers, username) VALUES (?, ?, ?, ?)",
            (chapter_id, quiz_part.strip(), answer_part.strip(), session['username'])
        )
        conn.commit()

        # Retrieve the inserted quiz ID using SQL Server's SCOPE_IDENTITY()
        # Retrieve last inserted quiz ID by filtering
        cursor.execute("""
                       SELECT TOP 1 quiz_id
                       FROM quizzes
                       WHERE chapter_id = ?
                         AND username = ?
                       ORDER BY date_created DESC
                       """, (chapter_id, session['username']))
        quiz_id_row = cursor.fetchone()

        if not quiz_id_row or not quiz_id_row[0]:
            conn.close()
            return "Failed to retrieve quiz ID.", 500

        quiz_id = quiz_id_row[0]
        conn.close()

        return redirect(url_for('solve_quiz', quiz_id=quiz_id))

    except Exception as e:
        return f"Quiz generation failed: {str(e)}", 500


@app.route('/solve_quiz/<int:quiz_id>', methods=['GET', 'POST'])
def solve_quiz(quiz_id):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Get quiz details
    cursor.execute("SELECT quiz_content, quiz_answers, chapter_id FROM quizzes WHERE quiz_id = ?", (quiz_id,))
    result = cursor.fetchone()
    if not result:
        conn.close()
        return "Quiz not found", 404

    quiz_text, answer_text, chapter_id = result

    # Get chapter name
    cursor.execute("SELECT chapter_name FROM chapters WHERE chapter_id = ?", (chapter_id,))
    chapter = cursor.fetchone()
    chapter_name = chapter[0] if chapter else "Unknown Chapter"

    conn.close()

    # Parse questions
    questions = []
    lines = quiz_text.strip().splitlines()
    current_q = {}

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line[0].isdigit() and ')' in line:
            if current_q:
                questions.append(current_q)
            current_q = {
                "question": line[line.find(')') + 1:].strip(),
                "options": []
            }
        elif line[0] in ['a', 'b', 'c'] and ')' in line:
            current_q["options"].append(line)
    if current_q:
        questions.append(current_q)

    # Parse answers
    correct_answers = {}
    answer_lines = answer_text.strip().splitlines()
    for line in answer_lines:
        if ')' in line:
            qnum, opt = line.split(')')
            correct_answers[int(qnum.strip())] = opt.strip()

    return render_template(
        "solve_quiz.html",
        questions=questions,
        quiz_id=quiz_id,
        chapter_name=chapter_name,
        correct_answers=correct_answers
    )

@app.route('/increment_quiz_count', methods=['POST'])
def increment_quiz_count():
    if 'username' not in session:
        return 'Unauthorized', 401

    username = session['username']
    try:
        connection = get_db_connection()
        cursor = connection.cursor()
        cursor.execute("UPDATE users SET quiz_completed = quiz_completed + 1 WHERE username = ?", (username,))
        connection.commit()
        return 'Success', 200
    except Exception as e:
        print("Error incrementing quiz count:", e)
        return 'Server Error', 500
    finally:
        connection.close()


@app.route('/quiz')
def quiz():
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT q.quiz_id, q.chapter_id, q.username, q.quiz_content, c.chapter_name, q.date_created
        FROM quizzes q
        JOIN chapters c ON q.chapter_id = c.chapter_id
        WHERE q.username = ?
        ORDER BY q.date_created DESC
    """, (session['username'],))

    quizzes = cursor.fetchall()
    conn.close()

    return render_template("quiz.html", quizzes=quizzes)



@app.route('/upload')
def upload():
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get all chapters for this user
    cursor.execute("""
        SELECT c.chapter_id, c.chapter_name, c.course_id
        FROM chapters c
        JOIN courses cr ON c.course_id = cr.course_id
        WHERE cr.username = ?
        AND c.chapter_id NOT IN (
            SELECT chapter_id FROM uploads WHERE username = ?
        )
    """, (session['username'], session['username']))
    chapters = cursor.fetchall()
    conn.close()

    return render_template("upload.html", chapters=chapters)

@app.route('/upload/<int:chapter_id>')
def upload_chapter_view(chapter_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT chapter_name, chapter_summary, course_id
        FROM chapters
        WHERE chapter_id = ?
    """, (chapter_id,))
    chapter = cursor.fetchone()
    conn.close()

    if not chapter:
        return "Chapter not found", 404

    return render_template("upload_summary.html", chapter_id=chapter_id,
                           chapter_name=chapter[0],
                           summary=chapter[1],
                           course_id=chapter[2])

@app.route('/upload_summary/<int:chapter_id>', methods=['POST'])
def upload_summary(chapter_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get course_id
    cursor.execute("SELECT course_id FROM chapters WHERE chapter_id = ?", (chapter_id,))
    course = cursor.fetchone()
    if not course:
        conn.close()
        return "Invalid chapter", 400

    # Insert into uploads
    cursor.execute("""
        INSERT INTO uploads (username, course_id, chapter_id, upload_date)
        VALUES (?, ?, ?, GETDATE())
    """, (session['username'], course[0], chapter_id))

    conn.commit()
    conn.close()

    return redirect(url_for('upload'))


@app.route('/load')
def load():
    if 'username' not in session:
        return redirect(url_for('login'))

    current_user = session['username']
    conn = get_db_connection()
    cursor = conn.cursor()

    query = request.args.get('query', '').strip()

    # Base query: all uploads not by the current user
    sql = """
        SELECT u.upload_id, u.username AS uploader, c.chapter_name
        FROM uploads u
        JOIN chapters c ON u.chapter_id = c.chapter_id
        WHERE u.username != ?
    """
    params = [current_user]

    # Add filter if search query is provided
    if query:
        sql += " AND (u.username LIKE ? OR c.chapter_name LIKE ?)"
        params.extend([f'%{query}%', f'%{query}%'])

    cursor.execute(sql, params)
    uploads = [dict(upload_id=row[0], uploader=row[1], chapter_name=row[2]) for row in cursor.fetchall()]
    conn.close()

    return render_template('load.html', uploads=uploads)


@app.route('/view_upload/<int:upload_id>')
def view_upload(upload_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT u.upload_id, u.chapter_id, u.course_id, u.username AS uploader, 
               ch.chapter_name, ch.chapter_summary
        FROM uploads u
        JOIN chapters ch ON u.chapter_id = ch.chapter_id
        WHERE u.upload_id = ?
    """, (upload_id,))
    chapter = cursor.fetchone()
    conn.close()

    if not chapter:
        return "Upload not found.", 404

    return render_template('view_upload.html', chapter=chapter)

@app.route('/load_chapter/<int:upload_id>', methods=['POST'])
def load_chapter(upload_id):
    if 'username' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get upload data
    cursor.execute("""
        SELECT u.username AS uploader, ch.chapter_name, ch.chapter_summary
        FROM uploads u
        JOIN chapters ch ON u.chapter_id = ch.chapter_id
        WHERE u.upload_id = ?
    """, (upload_id,))
    data = cursor.fetchone()

    if not data:
        conn.close()
        return "Upload not found", 404

    uploader = data[0]
    chapter_name = data[1]
    chapter_summary = data[2]

    # Check if course with uploader's name exists for current user
    cursor.execute("""
        SELECT course_id FROM courses
        WHERE course_name = ? AND username = ?
    """, (uploader, session['username']))
    course = cursor.fetchone()

    if course:
        course_id = course[0]
    else:
        # Create new course
        cursor.execute("""
            INSERT INTO courses (course_name, username)
            OUTPUT INSERTED.course_id
            VALUES (?, ?)
        """, (uploader, session['username']))
        course_id = cursor.fetchone()[0]

    # Insert chapter
    cursor.execute("""
        INSERT INTO chapters (chapter_name, course_id, chapter_summary)
        VALUES (?, ?, ?)
    """, (chapter_name, course_id, chapter_summary))

    conn.commit()
    conn.close()

    return redirect(url_for('load'))

if __name__ == '__main__':
    app.run(debug=True)